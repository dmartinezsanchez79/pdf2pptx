"""
Renderizador de presentaciones a formato .pptx.

Responsabilidad única: tomar una Presentation (modelo de dominio)
y escribir el archivo .pptx usando la plantilla de la universidad.

No genera contenido ni llama a la IA: solo transforma modelos en diapositivas.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation as PptxPresentation

from src.domain.models import Presentation, Slide, SlideType
from config.settings import TEMPLATE_PATH, LAYOUT_CONTENT


def render(presentation: Presentation, output_path: str | Path) -> Path:
    """
    Escribe la presentación en un archivo .pptx usando la plantilla universitaria.

    Returns:
        Path al archivo generado.
    """
    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)

    prs = PptxPresentation(TEMPLATE_PATH)

    # La plantilla trae 2 slides de ejemplo:
    #   slide 0 → portada real (text boxes TÍTULO / SUBTÍTULO + imagen)
    #   slide 1 → slide de contenido vacío (lo eliminamos, usamos el layout)
    _fill_cover(prs.slides[0], presentation.slides[0])
    _remove_slide(prs, 1)

    # Añadir el resto de slides (índice, desarrollo, conclusiones)
    for slide in presentation.slides[1:]:
        _add_content_slide(prs, slide)

    prs.save(output)
    return output


# ---------------------------------------------------------------------------
# Portada — modificar el slide existente en la plantilla
# ---------------------------------------------------------------------------

def _fill_cover(pptx_slide, slide: Slide) -> None:
    """
    Rellena la portada de la plantilla reemplazando los text boxes
    'TÍTULO' y 'SUBTÍTULO'.
    El subtítulo viene en bullets[0] del slide de portada.
    Preserva el formato original (fuente, tamaño, color) de cada text box.
    """
    subtitle = slide.bullets[0] if slide.bullets else "Presentación generada automáticamente"
    for shape in pptx_slide.shapes:
        if not shape.has_text_frame:
            continue
        current = shape.text_frame.text.strip().upper()
        if "TÍTULO" in current and "SUB" not in current:
            _replace_textbox(shape, slide.title)
        elif "SUBTÍTULO" in current or "SUB" in current:
            _replace_textbox(shape, subtitle)


def _replace_textbox(shape, new_text: str) -> None:
    """
    Sustituye el texto de un text box preservando el formato del primer run.
    Si el párrafo tiene runs, reutiliza el primero (conserva fuente/color/tamaño).
    """
    tf = shape.text_frame
    para = tf.paragraphs[0]
    if para.runs:
        para.runs[0].text = new_text
        # Eliminar runs adicionales si los hubiera
        for run in para.runs[1:]:
            run.text = ""
    else:
        para.add_run().text = new_text


# ---------------------------------------------------------------------------
# Slides de contenido — añadir desde el layout
# ---------------------------------------------------------------------------

def _add_content_slide(prs: PptxPresentation, slide: Slide) -> None:
    """
    Añade una diapositiva de contenido (índice, desarrollo o conclusiones)
    usando el layout TITLE_AND_BODY de la plantilla.
    Los bullets se escriben párrafo a párrafo para que el template
    aplique su propio estilo de viñeta sin duplicarlo.
    """
    layout = prs.slide_layouts[LAYOUT_CONTENT]
    pptx_slide = prs.slides.add_slide(layout)

    _set_placeholder(pptx_slide, idx=0, text=slide.title)

    if slide.bullets:
        _set_body_bullets(pptx_slide, slide.bullets)


def _set_body_bullets(pptx_slide, bullets: list[str]) -> None:
    """
    Escribe los bullets en el placeholder de cuerpo (idx=1)
    usando un párrafo por bullet.
    No añade prefijo manual — el estilo de viñeta viene del layout.
    """
    for ph in pptx_slide.placeholders:
        if ph.placeholder_format.idx != 1:
            continue
        tf = ph.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = bullets[0]
        tf.paragraphs[0].level = 0
        for bullet in bullets[1:]:
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 0
        return


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _set_placeholder(slide, idx: int, text: str) -> None:
    """Asigna texto a un placeholder por índice. No falla si no existe."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            ph.text = text
            return


def _remove_slide(prs: PptxPresentation, index: int) -> None:
    """Elimina un slide de la presentación por su índice."""
    sldIdLst = prs.slides._sldIdLst
    sldId = sldIdLst[index]
    rId = sldId.get(
        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
    )
    prs.part.drop_rel(rId)
    sldIdLst.remove(sldId)
